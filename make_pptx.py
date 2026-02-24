from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

W = Inches(13.33)
H = Inches(7.5)

def rgb(r, g, b):
    return RGBColor(r, g, b)

def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    # 배경 흰색
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(255, 255, 255)
    return slide

def box(slide, left, top, width, height, bg=None, border=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background() if border is None else None
    if bg:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(*bg)
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = rgb(*border)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, left, top, width, height,
        size=14, bold=False, color=(30,30,30),
        align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(*color)
    return txBox

def add_rounded_box(slide, left, top, width, height, bg, radius_pt=6):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.adjustments[0] = 0.05
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(*bg)
    shape.line.fill.background()
    return shape

# ─────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ═══════════════════════════════════════
# 슬라이드 1 — 타이틀
# ═══════════════════════════════════════
slide = add_slide(prs)
# 배경 그라데이션 느낌 — 상단 인디고 밴드
b = box(slide, 0, 0, 13.33, 7.5, bg=(67, 56, 202))
txt(slide, "실험 통합 DB 필요성", 0.6, 2.2, 12, 1.2,
    size=40, bold=True, color=(255,255,255), align=PP_ALIGN.CENTER)
txt(slide, "세 가지 DB를 엮었을 때 생기는 맥락  —  단절 전 / 후 비교",
    0.6, 3.5, 12, 0.7, size=18, color=(199,210,254), align=PP_ALIGN.CENTER)
txt(slide, "📁 과제(Projects)  →  🧪 실험조건(Experiments)  →  📋 Split Table  →  📈 Lot Summary",
    0.6, 4.4, 12, 0.6, size=13, color=(165,180,252), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# 슬라이드 2 — DB 연결 구조
# ═══════════════════════════════════════
slide = add_slide(prs)

# 헤더 바
b = box(slide, 0, 0, 13.33, 0.9, bg=(67, 56, 202))
txt(slide, "세 DB가 만드는 체인 구조", 0.4, 0.12, 12, 0.66,
    size=20, bold=True, color=(255,255,255))

CARDS = [
    ((67,56,202),  "📁", "과제\n(Projects)",          "• 과제명 / 목표\n• 담당자 / 일정\n• 과제 분류",      "iacpj_nm (PK)"),
    ((109,40,217), "🧪", "실험 조건\n(Experiments)",  "• Plan ID / 평가항목\n• 평가공정 / LOT 코드\n• 실험 목표", "plan_id (PK)\niacpj_nm (FK)"),
    ((37,99,235),  "📋", "Split Table",               "• OPER_ID / 공정명\n• 장비 / Recipe / 변수값\n• WF 배정", "plan_id (FK)\noper_id"),
    ((5,150,105),  "📈", "실험 결과\n(Lot Summary)",  "• LOT ID / 측정값\n• Inline / Outline 결과\n• Pass/Fail 판정", "lot_id\nplan_id (FK)"),
]

x_starts = [0.3, 3.45, 6.6, 9.75]
card_w = 2.9
card_h = 5.3
top = 1.1

for i, (col, icon, title, desc, key) in enumerate(CARDS):
    x = x_starts[i]
    b = box(slide, x, top, card_w, card_h, bg=col)
    txt(slide, icon, x+0.15, top+0.15, card_w-0.3, 0.55, size=28)
    txt(slide, title, x+0.15, top+0.75, card_w-0.3, 0.8,
        size=14, bold=True, color=(255,255,255))
    txt(slide, desc, x+0.15, top+1.6, card_w-0.3, 2.2,
        size=11, color=(210,230,255))
    # key 배지
    b3 = box(slide, x+0.15, top+4.55, card_w-0.3, 0.55, bg=(30,20,80))
    txt(slide, key, x+0.2, top+4.6, card_w-0.4, 0.45,
        size=9, color=(180,200,255))

    # 화살표 (마지막 카드 제외)
    if i < 3:
        arr_x = x_starts[i] + card_w + 0.05
        box(slide, arr_x, top + card_h/2 - 0.03, 0.35, 0.06, bg=(150,150,220))
        # 화살표 머리
        arr = slide.shapes.add_shape(5, Inches(arr_x+0.28), Inches(top+card_h/2-0.13),
                                     Inches(0.17), Inches(0.26))
        arr.fill.solid(); arr.fill.fore_color.rgb = rgb(150,150,220)
        arr.line.fill.background()

# 1:N 레이블
for i in range(3):
    ax = x_starts[i] + card_w + 0.12
    txt(slide, "1:N", ax, top + card_h/2 - 0.32, 0.25, 0.25,
        size=9, color=(120,120,180), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# 슬라이드 3 — 단절 상태의 한계
# ═══════════════════════════════════════
slide = add_slide(prs)
b = box(slide, 0, 0, 13.33, 0.9, bg=(220, 38, 38))
txt(slide, "⚠️  통합 DB 없이 — 단절 상태의 제약", 0.4, 0.12, 12, 0.66,
    size=20, bold=True, color=(255,255,255))

LIMITS = [
    ("📁", "과제 정보 고립",       "과제 목표·담당자·일정이 그룹웨어·문서에 분산 저장.\n실험 조건과 연결 고리 없음.",              "\"이 실험이 어느 과제 소속?\" → 담당자에게 직접 문의"),
    ("📊", "Split 조건 개인 관리", "Split Table이 담당자 엑셀에만 존재.\n공유·버전 관리 불가.",                                    "담당자 부재 시 조건 확인 불가, 이력 소실"),
    ("🔬", "LOT-조건 연결 단절",   "어떤 LOT이 어떤 Split으로 진행됐는지\nMES와 실험 계획이 별도 시스템.",                        "\"이 LOT 결과가 왜 이렇게?\" → 수동 대조 작업"),
    ("🔍", "유사 실험 탐색 불가",  "과거 실험을 찾으려면 담당자 기억 또는\n수백 개 엑셀 파일 수동 탐색.",                          "검색에 수 시간 소요, 관련 실험을 놓치는 경우 다반사"),
    ("🤖", "AI 활용 불가",         "데이터가 분절되어 AI가 맥락을 파악 못함.\n단순 텍스트 검색만 가능.",                           "\"이 과제 최적 조건이 뭔지\" AI에 물어봐도 답 없음"),
    ("👤", "지식이 사람에게 귀속", "실험 노하우·이력이 개인 파일·기억에 의존.\n담당자 교체 시 처음부터 다시.",                     "신입 투입 시 온보딩에 수 주 소요, 실수 반복"),
]

cols = [(0.25, 4.3), (6.85, 4.3)]
rows_per_col = 3

for idx, (icon, title, desc, pain) in enumerate(LIMITS):
    col = idx // rows_per_col
    row = idx % rows_per_col
    x = cols[col][0]
    y = 1.05 + row * 2.1

    b = box(slide, x, y, 6.3, 1.9, bg=(254,242,242), border=(252,165,165))
    txt(slide, icon, x+0.1, y+0.1, 0.5, 0.5, size=22)
    txt(slide, title, x+0.65, y+0.1, 5.5, 0.35, size=13, bold=True, color=(153,27,27))
    txt(slide, desc,  x+0.65, y+0.5, 5.5, 0.7,  size=10, color=(100,20,20))
    # pain 배지
    b2 = box(slide, x+0.1, y+1.35, 6.05, 0.42, bg=(254,202,202))
    txt(slide, "💢 " + pain, x+0.2, y+1.38, 5.9, 0.36,
        size=9, bold=True, color=(185,28,28))

# ═══════════════════════════════════════
# 슬라이드 4 — 통합 후 가능한 탐색 (1)
# ═══════════════════════════════════════
slide = add_slide(prs)
b = box(slide, 0, 0, 13.33, 0.9, bg=(5, 150, 105))
txt(slide, "✅  통합 DB 구축 후 — 가능해지는 탐색 (1/2)", 0.4, 0.12, 12, 0.66,
    size=20, bold=True, color=(255,255,255))

AFTER1 = [
    ((67,56,202), (238,242,255), (165,180,252), "과제 → 실험",
     "\"ESL CMP 신뢰성 향상 과제에서 진행된 실험 전체 목록은?\"",
     "Projects → Experiments JOIN → Plan ID별 평가항목·공정·LOT 즉시 반환",
     "과제명 하나로 소속 실험 전체 + Split 조건 + 결과까지 원클릭 조회"),
    ((109,40,217),(245,243,255),(196,181,253), "실험 → Split",
     "\"PLN-2025-003 실험의 공정별 Split 조건과 배정 WF를 보여줘\"",
     "Experiments → Split Table JOIN → OPER_ID별 조건·장비·Recipe·WF 매핑 반환",
     "base/s1/s2/s3 각 Split이 어떤 WF에서 어떤 조건으로 진행됐는지 한눈에"),
    ((37,99,235), (239,246,255),(147,197,253), "조건 → 결과",
     "\"Slurry flow -15% 조건으로 진행한 LOT의 CMP 균일도 결과는?\"",
     "Split Table(work_cond_desc 검색) → plan_id → Lot Summary 결과값 추적",
     "조건 텍스트 하나로 해당 조건이 적용된 LOT·결과를 역방향 추적"),
]

for idx, (hcol, bgcol, badgecol, badge, q, a, detail) in enumerate(AFTER1):
    y = 1.05 + idx * 2.1
    b = box(slide, 0.25, y, 12.83, 1.9, bg=bgcol, border=badgecol)
    # 배지
    b2 = box(slide, 0.35, y+0.12, 1.3, 0.32, bg=badgecol)
    txt(slide, badge, 0.38, y+0.13, 1.24, 0.28,
        size=9, bold=True, color=hcol, align=PP_ALIGN.CENTER)
    txt(slide, "🗨️ " + q, 0.35, y+0.52, 12.5, 0.42,
        size=12, bold=True, color=hcol)
    txt(slide, "▸ " + a, 0.35, y+0.98, 12.5, 0.38, size=10, color=(80,80,80))
    txt(slide, "→ " + detail, 0.35, y+1.42, 12.5, 0.35,
        size=10, bold=True, color=hcol)

# ═══════════════════════════════════════
# 슬라이드 5 — 통합 후 가능한 탐색 (2)
# ═══════════════════════════════════════
slide = add_slide(prs)
b = box(slide, 0, 0, 13.33, 0.9, bg=(5, 150, 105))
txt(slide, "✅  통합 DB 구축 후 — 가능해지는 탐색 (2/2)", 0.4, 0.12, 12, 0.66,
    size=20, bold=True, color=(255,255,255))

AFTER2 = [
    ((5,150,105),  (236,253,245),(110,231,183), "유사 실험 탐색",
     "\"CMP 균일도 평가 실험 중 WIW 2% 이하 달성한 과거 조건 찾아줘\"",
     "eval_item LIKE 검색 → Split 조건 열람 → Lot Summary 결과 필터링",
     "평가항목·공정 키워드 → 조건 → 결과까지 한 번에 탐색, AI 추천 가능"),
    ((180,83,9),   (255,251,235),(253,230,138), "이력 추적",
     "\"OPER_ID CMP0100 공정에서 지난 1년간 사용된 모든 Recipe 이력은?\"",
     "Split Table WHERE oper_id='CMP0100' → recipe_id 집계 → 시계열 정렬",
     "공정 단위 Recipe 변경 이력 추적, 이상 발생 시 원인 조건 즉시 특정"),
    ((159,18,57),  (255,241,242),(253,164,175), "AI 맥락 제공",
     "\"신규 CMP 실험 설계할 때 과거 성공 조건 기반으로 추천해줘\"",
     "Projects + Experiments + Split + Lot Summary 전체 맥락 → AI 종합 분석",
     "단절 DB에서 불가능했던 맥락 기반 AI 추천이 세 DB 연결로 실현"),
]

for idx, (hcol, bgcol, badgecol, badge, q, a, detail) in enumerate(AFTER2):
    y = 1.05 + idx * 2.1
    b = box(slide, 0.25, y, 12.83, 1.9, bg=bgcol, border=badgecol)
    b2 = box(slide, 0.35, y+0.12, 1.5, 0.32, bg=badgecol)
    txt(slide, badge, 0.38, y+0.13, 1.44, 0.28,
        size=9, bold=True, color=hcol, align=PP_ALIGN.CENTER)
    txt(slide, "🗨️ " + q, 0.35, y+0.52, 12.5, 0.42,
        size=12, bold=True, color=hcol)
    txt(slide, "▸ " + a, 0.35, y+0.98, 12.5, 0.38, size=10, color=(80,80,80))
    txt(slide, "→ " + detail, 0.35, y+1.42, 12.5, 0.35,
        size=10, bold=True, color=hcol)

# ═══════════════════════════════════════
# 슬라이드 6 — 결론
# ═══════════════════════════════════════
slide = add_slide(prs)
b = box(slide, 0, 0, 13.33, 7.5, bg=(67, 56, 202))

txt(slide, "🔗", 5.9, 1.0, 1.5, 1.2, size=60, align=PP_ALIGN.CENTER)
txt(slide, "단순 저장이 아닌  \"맥락 있는 DB\"", 0.5, 2.3, 12.33, 1.0,
    size=32, bold=True, color=(255,255,255), align=PP_ALIGN.CENTER)

summary = (
    "과제  →  실험 조건  →  Split  →  결과로 이어지는 체인이 완성될 때,\n"
    "단순 데이터 조회를 넘어  AI가 맥락을 이해하고 새 실험을 설계하는 지식 기반이 됩니다.\n\n"
    "담당자가 바뀌어도 지식은 DB에 남고,  다음 실험은 더 빠르고 정확하게 설계됩니다."
)
txt(slide, summary, 1.0, 3.5, 11.33, 2.2,
    size=15, color=(199,210,254), align=PP_ALIGN.CENTER)

# 하단 3개 포인트
POINTS = [
    ("🕐", "수 시간 → 즉시", "탐색 시간"),
    ("📚", "기억 → DB", "지식 귀속"),
    ("🤖", "불가 → 가능", "AI 활용"),
]
for i, (icon, val, label) in enumerate(POINTS):
    x = 2.0 + i * 3.5
    b2 = box(slide, x, 5.8, 3.0, 1.3, bg=(79,70,229))
    txt(slide, icon, x+0.1, 5.85, 0.7, 0.6, size=24)
    txt(slide, val,   x+0.8, 5.88, 2.0, 0.5, size=14, bold=True, color=(255,255,255))
    txt(slide, label, x+0.8, 6.35, 2.0, 0.4, size=10, color=(165,180,252))

out_path = "/Users/dail/Downloads/통합DB_필요성.pptx"
prs.save(out_path)
print(f"저장 완료: {out_path}")

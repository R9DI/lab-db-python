from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W = Inches(13.33)
H = Inches(7.5)

def rgb(r, g, b): return RGBColor(r, g, b)

def box(slide, left, top, width, height, bg=None, border_col=None, border_pt=1):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    if bg:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(*bg)
    else:
        shape.fill.background()
    if border_col:
        shape.line.color.rgb = rgb(*border_col)
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, left, top, width, height,
        size=12, bold=False, color=(30,30,30), align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(*color)
    return txBox

def multiline_txt(slide, lines, left, top, width, height,
                  size=10, bold=False, color=(80,80,80), spacing_pt=2):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(spacing_pt)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(*color)

# ── 프레젠테이션 설정
prs = Presentation()
prs.slide_width = W
prs.slide_height = H

slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = rgb(248, 249, 252)

# ══════════════════════════════════════════════════
# 1. 헤더 바
# ══════════════════════════════════════════════════
box(slide, 0, 0, 13.33, 0.68, bg=(67, 56, 202))
txt(slide, "통합 DB 필요성", 0.35, 0.06, 7, 0.38,
    size=22, bold=True, color=(255,255,255))
txt(slide, "데이터가 분산될 때 발생하는 4가지 문제의 연쇄  —  그리고 통합이 만드는 변화",
    0.35, 0.38, 9, 0.26, size=10, color=(199,210,254))

# ══════════════════════════════════════════════════
# 2. DB 체인 바
# ══════════════════════════════════════════════════
box(slide, 0, 0.72, 13.33, 0.52, bg=(79, 70, 229))
txt(slide, "🗄️  세 DB가 만드는 체인:", 0.3, 0.79, 2.8, 0.32,
    size=10, bold=True, color=(255,255,255))

CHAIN = ["📁 Projects", "🧪 Experiments", "📋 Split Table", "📈 Lot Summary"]
for i, label in enumerate(CHAIN):
    cx = 3.2 + i * 2.55
    box(slide, cx, 0.79, 2.1, 0.3, bg=(99,91,240), border_col=(165,180,252))
    txt(slide, label, cx+0.08, 0.8, 1.94, 0.28,
        size=9, bold=True, color=(255,255,255), align=PP_ALIGN.CENTER)
    if i < 3:
        txt(slide, "→", cx+2.13, 0.8, 0.3, 0.28,
            size=11, color=(199,210,254), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════
# 3. 섹션 레이블
# ══════════════════════════════════════════════════
txt(slide, "▸  분리 저장이 만드는 4가지 문제의 연쇄",
    0.3, 1.3, 8, 0.25, size=9, bold=True, color=(120,120,160))

# ══════════════════════════════════════════════════
# 4. 4개 문제 카드
# ══════════════════════════════════════════════════
PROBLEMS = [
    {
        "num": "01", "icon": "🗂️", "title": "Data 저장 비효율",
        "subtitle": "분산 저장 → 관리 불가",
        "hdr": (239,68,68), "bg": (254,242,242), "border": (252,165,165),
        "nbg": (220,38,38), "pain_bg": (254,202,202), "pain_c": (185,28,28),
        "desc": "과제·실험·Split·결과가 각기 다른 시스템·엑셀·문서에 저장됩니다. 버전이 엇갈리며 최신 파일이 어딘지 알 수 없습니다.",
        "items": ["과제 정보 — 그룹웨어 문서", "실험 조건 — 개인 엑셀", "Split 조건 — 별도 공정 파일", "측정 결과 — MES/분석 보고서"],
        "pain": "같은 실험 하나를 파악하려면 4곳을 뒤져야 한다",
    },
    {
        "num": "02", "icon": "🔗", "title": "맥락 부재",
        "subtitle": "분리 저장 → 실험도 맥락 없음",
        "hdr": (249,115,22), "bg": (255,247,237), "border": (253,186,116),
        "nbg": (234,88,12), "pain_bg": (254,215,170), "pain_c": (154,52,18),
        "desc": "동일한 실험임에도 'Split 조건이 어떤 과제의 어떤 목표를 위해 설계됐는지' 알 수 없습니다. 결과만 있고 배경이 없는 상태.",
        "items": ["\"이 LOT 조건이 뭐였지?\" → 담당자 문의", "\"어떤 목표를 검증하려 했지?\" → 문서 탐색", "조건·결과·목표가 각각 단절된 섬", "담당자 퇴사 시 배경 지식 완전 소실"],
        "pain": "결과 숫자는 있지만 왜 그 조건인지 아무도 모른다",
    },
    {
        "num": "03", "icon": "🔍", "title": "연속적 검색 불가",
        "subtitle": "맥락 없음 → 이어서 탐색 불가",
        "hdr": (245,158,11), "bg": (255,251,235), "border": (253,230,138),
        "nbg": (217,119,6), "pain_bg": (254,243,199), "pain_c": (146,64,14),
        "desc": "\"이 조건으로 진행한 LOT의 결과는?\" 같은 질문도 여러 시스템을 수동 대조해야 합니다. 검색이 깊어질수록 시간은 기하급수적 증가.",
        "items": ["과제명으로 실험 찾기 → 수동 엑셀 필터", "실험 조건으로 결과 추적 → 수동 대조", "유사 실험 탐색 → 기억 또는 전체 파일 탐색", "이력 추적 → 구두 확인 필요"],
        "pain": "한 번 검색에 수 시간, 그래도 빠뜨리는 케이스 발생",
    },
    {
        "num": "04", "icon": "🤖", "title": "AI 활용 불가",
        "subtitle": "맥락 없음 → AI도 답을 모른다",
        "hdr": (124,58,237), "bg": (245,243,255), "border": (196,181,253),
        "nbg": (109,40,217), "pain_bg": (221,214,254), "pain_c": (76,29,149),
        "desc": "AI는 구조화된 맥락이 있어야 추론합니다. 데이터가 분절되면 AI에게 최적 조건을 물어봐도 참조할 연결 정보가 없어 의미있는 답을 줄 수 없습니다.",
        "items": ["과거 성공 조건 기반 추천 → 불가 (단절)", "유사 실험 자동 제안 → 불가 (맥락 없음)", "조건 변경 영향 예측 → 불가 (결과 미연결)", "신규 실험 설계 보조 → 불가 (배경 없음)"],
        "pain": "데이터는 쌓이는데, AI는 그걸 활용할 수 없다",
    },
]

card_w = 3.12
card_top = 1.6
card_h = 4.35
gap = 0.09

for i, p in enumerate(PROBLEMS):
    x = 0.18 + i * (card_w + gap)

    # 카드 배경
    box(slide, x, card_top, card_w, card_h, bg=p["bg"], border_col=p["border"])

    # 헤더
    box(slide, x, card_top, card_w, 0.72, bg=p["hdr"])

    # 번호
    txt(slide, p["num"], x+0.1, card_top+0.02, 0.65, 0.65,
        size=28, bold=True, color=(*p["nbg"], ), align=PP_ALIGN.LEFT)
    # 실제 번호를 밝게
    txt(slide, p["num"], x+0.1, card_top+0.02, 0.65, 0.65,
        size=28, bold=True, color=(255,255,255), align=PP_ALIGN.LEFT)

    txt(slide, p["icon"] + " " + p["title"],
        x+0.75, card_top+0.06, card_w-0.85, 0.32,
        size=11, bold=True, color=(255,255,255))
    txt(slide, p["subtitle"],
        x+0.75, card_top+0.38, card_w-0.85, 0.26,
        size=8.5, color=(255,220,180))

    # 설명 텍스트
    txt(slide, p["desc"],
        x+0.15, card_top+0.78, card_w-0.3, 0.95,
        size=8.5, color=(70,70,90))

    # 구분선
    box(slide, x+0.15, card_top+1.78, card_w-0.3, 0.01, bg=(200,200,220))

    # 불릿 항목
    multiline_txt(slide, ["• " + it for it in p["items"]],
                  x+0.15, card_top+1.84, card_w-0.3, 1.5,
                  size=8.5, color=(80,80,100), spacing_pt=1.5)

    # Pain 배지
    box(slide, x+0.12, card_top+3.72, card_w-0.24, 0.48, bg=p["pain_bg"])
    txt(slide, "💢 " + p["pain"],
        x+0.18, card_top+3.76, card_w-0.36, 0.4,
        size=8, bold=True, color=p["pain_c"])

    # 화살표 (마지막 제외)
    if i < 3:
        ax = x + card_w + 0.01
        ay = card_top + card_h/2
        box(slide, ax, ay-0.02, gap+0.01, 0.04, bg=(150,150,200))

# ══════════════════════════════════════════════════
# 5. 연쇄 레이블
# ══════════════════════════════════════════════════
CHAIN_LABELS = ["저장 비효율", "→  맥락 부재", "→  연속 검색 불가", "→  AI 활용 불가"]
CHAIN_COLORS = [(220,38,38), (234,88,12), (180,100,0), (109,40,217)]
lw = 3.12
for i, (label, col) in enumerate(zip(CHAIN_LABELS, CHAIN_COLORS)):
    lx = 0.18 + i * (lw + gap)
    txt(slide, label, lx, card_top + card_h + 0.06, lw, 0.22,
        size=8.5, bold=True, color=col, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════
# 6. After 바 (하단)
# ══════════════════════════════════════════════════
after_top = card_top + card_h + 0.32
after_h = 7.5 - after_top - 0.08
box(slide, 0, after_top, 13.33, after_h, bg=(67,56,202))

txt(slide, "✅  통합 DB 구축 후 — 4가지 문제 동시 해결",
    0.3, after_top+0.1, 7, 0.3, size=11, bold=True, color=(255,255,255))

AFTER = [
    ("🗄️", "단일 저장소", "과제·실험·Split·결과\n하나의 DB에 연결 저장"),
    ("🔗", "완전한 맥락", "실험 하나가 목표부터\n결과까지 이어진 스토리"),
    ("🔍", "연속 탐색", "어느 방향으로도\n즉시 검색·추적 가능"),
    ("🤖", "AI 추천 가능", "맥락 기반 유사 실험 제안\n신규 실험 설계 보조"),
]
aw = 3.0
for i, (icon, title, desc) in enumerate(AFTER):
    ax = 0.25 + i * (aw + 0.12)
    box(slide, ax, after_top+0.46, aw, after_h-0.54, bg=(88,80,236))
    txt(slide, icon + "  " + title,
        ax+0.15, after_top+0.52, aw-0.2, 0.3,
        size=10, bold=True, color=(255,255,255))
    txt(slide, desc,
        ax+0.15, after_top+0.82, aw-0.2, after_h-1.0,
        size=8.5, color=(199,210,254))

# 마지막 문구
txt(slide,
    "과제→실험→Split→결과 체인 완성 시,  AI가 맥락을 이해하고 새 실험을 설계하는 지식 기반이 됩니다.",
    0.25, after_top+0.12, 13.0, 0.28,
    size=8.5, color=(165,180,252), align=PP_ALIGN.RIGHT)

out_path = "/Users/dail/Downloads/통합DB_필요성_1page.pptx"
prs.save(out_path)
print(f"저장 완료: {out_path}")
